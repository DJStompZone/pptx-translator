import subprocess
import json
from pptx import Presentation

def extract_notes_to_json(ppt_path, json_path):
    ppt = Presentation(ppt_path)
    notes = [(slide_id, slide.notes_slide.notes_text_frame.text) for slide_id, slide in enumerate(ppt.slides)]
    with open(json_path, 'w', encoding='utf-8') as fp: 
        json.dump(notes, fp)

def update_presentation_with_translated_notes(ppt_path, json_translated_path):
    ppt = Presentation(ppt_path)
    with open(json_translated_path, 'r', encoding='utf-8') as fp: 
        translated_notes = json.load(fp)
    
    for item in translated_notes:
        slide_id = item['original'][0]
        translated_text = item['translated']
        ppt.slides[slide_id].notes_slide.notes_text_frame.text = translated_text

    ppt.save(ppt_path)

def run_translation_script():
    subprocess.run(["node", "translate.js"], check=True)

def main():
    ppt_path = '../session10.pptx'
    json_path = 'output.json'
    json_translated_path = 'translated.json'
    
    extract_notes_to_json(ppt_path, json_path)
    
    run_translation_script()
    
    update_presentation_with_translated_notes(ppt_path, json_translated_path)

if __name__ == "__main__":
    main()
