import subprocess
import json
from pptx import Presentation
from pptx.util import Inches

def extract_text_to_json(ppt_path, json_path):
    ppt = Presentation(ppt_path)
    text_boxes = []
    for slide_id, slide in enumerate(ppt.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.text_frame.text:
                text_boxes.append((slide_id, shape.text_frame.text))
    
    with open(json_path, 'w', encoding='utf-8') as fp:
        json.dump(text_boxes, fp)

def update_presentation_with_translated_text(ppt_path, json_translated_path):
    ppt = Presentation(ppt_path)
    with open(json_translated_path, 'r', encoding='utf-8') as fp:
        translated_texts = json.load(fp)

    for item in translated_texts:
        slide_id = item['original'][0]
        translated_text = item['translated']
        slide = ppt.slides[slide_id]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text == item['original'][1]:
                shape.text_frame.text = translated_text

    ppt.save(ppt_path)

def run_translation_script():
    subprocess.run(["node", "translate_boxes.js"], check=True)

def main():
    ppt_path = '../session10.pptx'
    json_path = 'output_boxes.json'
    json_translated_path = 'translated_boxes.json'
    
    extract_text_to_json(ppt_path, json_path)
    
    run_translation_script()
    
    update_presentation_with_translated_text(ppt_path, json_translated_path)

if __name__ == "__main__":
    main()
