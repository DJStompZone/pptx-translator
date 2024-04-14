import asyncio
import json

from local_io import extract_to_json
from local_types import ExtractionMode
from pptx import Presentation
from translate import Config, SlideElement, TranslationResult, translate_slide


def update_presentation(ppt_path, mode, translated_texts):
    """
    Update the given PowerPoint presentation with translated texts.

    Args:
        ppt_path (str): The file path of the PowerPoint presentation.
        mode (ExtractionMode): The mode of extraction for translated texts.
        translated_texts (list): A list of translated texts.

    Returns:
        None
    """
    ppt = Presentation(ppt_path)
    for response_item in translated_texts:
        slide_id = response_item.original.slide_id
        original_text = response_item.original.text
        translated_text = response_item.translated

        if not response_item.error and translated_text:
            slide = ppt.slides[slide_id]
            if mode == ExtractionMode.TEXT_BOXES:
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text == original_text:
                        shape.text_frame.text = translated_text
            elif mode == ExtractionMode.NOTES:
                if slide.has_notes_slide:
                    slide.notes_slide.notes_text_frame.text = translated_text

    ppt.save(ppt_path)


async def translation_coordinator(ppt_path: str, json_path: str, json_translated_path: str, mode: ExtractionMode):
    """
    Coordinates the translation process for a PowerPoint presentation.

    Args:
        ppt_path (str): The path to the PowerPoint presentation file.
        json_path (str): The path to the JSON file used for extraction and translation.
        json_translated_path (str): The path to the JSON file where translated texts will be saved.
        mode (ExtractionMode): The extraction mode for extracting elements from the presentation.

    Returns:
        None
    """
    extract_to_json(ppt_path, json_path, mode)

    with open(json_path, 'r', encoding='utf-8') as fp:
        elements_dicts = json.load(fp)
    
    translated_texts: list[TranslationResult] = []
    config = Config.load_from_file('../config.json', json_path, json_translated_path)
    for element_dict in elements_dicts:
        element = SlideElement(**element_dict)
        result = await translate_slide(config, element)
        translated_texts.append(result)

    with open(json_translated_path, 'w', encoding='utf-8') as fp:
        json.dump(translated_texts, fp)

    # Update presentation with translated texts
    update_presentation(ppt_path, mode, translated_texts)

def main(ppt_path: str = '../session10.pptx', json_path_text: str = 'output_boxes.json', json_translated_path_text: str = 'translated_boxes.json', json_path_notes: str = 'output.json', json_translated_path_notes: str = 'translated.json'):
    asyncio.run(translation_coordinator(ppt_path, json_path_text, json_translated_path_text, mode=ExtractionMode.TEXT_BOXES))
    asyncio.run(translation_coordinator(ppt_path, json_path_notes, json_translated_path_notes, mode=ExtractionMode.NOTES))

if __name__ == "__main__":
    main()
