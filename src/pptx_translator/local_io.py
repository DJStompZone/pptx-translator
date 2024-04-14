import json

from local_types import ExtractionMode, SlideElement
from pptx import Presentation


def extract_to_json(ppt_path: str, json_path: str, mode: ExtractionMode):
    """Extracts text from text boxes or notes depending on the mode."""
    ppt = Presentation(ppt_path)
    data: list[SlideElement] = []
    for slide_id, slide in enumerate(ppt.slides):
        if mode == ExtractionMode.TEXT_BOXES:
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text:
                    data.append(SlideElement(slide_id=slide_id, text=shape.text_frame.text))
        elif mode == ExtractionMode.NOTES:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                data.append(SlideElement(slide_id=slide_id, text=slide.notes_slide.notes_text_frame.text))

    with open(json_path, 'w', encoding='utf-8') as fp:
        json.dump([element.__dict__ for element in data], fp)